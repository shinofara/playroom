"""Cotomo 1分ピッチ PPTX生成スクリプト"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- 定数 ---
SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)
IMG = os.path.join(os.path.dirname(__file__), "images")

# ブランドカラー
BG_LIGHT = RGBColor(0xF3, 0xFB, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY = RGBColor(0x48, 0x48, 0x48)
TEXT_SECONDARY = RGBColor(0x6B, 0x6B, 0x6B)
TEXT_DARK = RGBColor(0x57, 0x57, 0x57)
BRAND_GREEN = RGBColor(0x40, 0xB2, 0x87)
BUBBLE_GREEN = RGBColor(0x6C, 0xC7, 0xA5)
BUBBLE_YELLOW = RGBColor(0xD1, 0xC4, 0x5C)
BUBBLE_PINK = RGBColor(0xDA, 0x76, 0xD7)
ACCENT_PINK = RGBColor(0xDD, 0x82, 0xDA)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


def set_bg_color(slide, color):
    """スライド背景色を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bubble(slide, left, top, width, height, color, opacity=0.35):
    """装飾バブル（半透明楕円）を追加"""
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    # XML直接操作で透明度設定
    sp_pr = shape._element.spPr
    solid_fill = sp_pr.find(f'{{{ns}}}solidFill')
    if solid_fill is not None:
        srgb = solid_fill.find(f'{{{ns}}}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, f'{{{ns}}}alpha')
            alpha_elem.set('val', str(int(opacity * 100000)))
    shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=32,
                 bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER,
                 font_name='Noto Sans JP', italic=False):
    """テキストボックスを追加"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txbox


def add_rich_text(slide, left, top, width, height, runs, alignment=PP_ALIGN.CENTER):
    """複数のrunを持つリッチテキストボックスを追加"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, run_def in enumerate(runs):
        if i == 0:
            r = p.runs[0] if p.runs else p.add_run()
        else:
            r = p.add_run()
        r.text = run_def.get('text', '')
        r.font.size = Pt(run_def.get('size', 32))
        r.font.bold = run_def.get('bold', False)
        r.font.italic = run_def.get('italic', False)
        r.font.color.rgb = run_def.get('color', TEXT_PRIMARY)
        r.font.name = run_def.get('font', 'Noto Sans JP')
    return txbox


def add_multiline_text(slide, left, top, width, height, lines, alignment=PP_ALIGN.CENTER):
    """複数行テキストボックスを追加"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line_def in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_def.get('text', '')
        p.font.size = Pt(line_def.get('size', 32))
        p.font.bold = line_def.get('bold', False)
        p.font.color.rgb = line_def.get('color', TEXT_PRIMARY)
        p.font.name = line_def.get('font', 'Noto Sans JP')
        p.alignment = alignment
        p.space_after = Pt(line_def.get('space_after', 6))
    return txbox


def add_image_safe(slide, path, left, top, width=None, height=None):
    """画像を安全に追加（ファイル存在チェック付き）"""
    full = os.path.join(IMG, path)
    if not os.path.exists(full):
        print(f"  [WARN] 画像なし: {full}")
        return None
    kwargs = {'left': left, 'top': top}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    return slide.shapes.add_picture(full, **kwargs)


def add_rounded_rect(slide, left, top, width, height, fill_color=WHITE, line=False):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_note(slide, text):
    """スピーカーノートを追加"""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# =================================================================
# S1: タイトル
# =================================================================
s1 = prs.slides.add_slide(blank_layout)
set_bg_color(s1, BG_LIGHT)
add_image_safe(s1, 'bg-characters.png', Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_bubble(s1, Inches(-0.5), Inches(-0.8), Inches(3), Inches(3), BUBBLE_GREEN)
add_bubble(s1, Inches(10), Inches(5), Inches(2.5), Inches(2.5), BUBBLE_PINK)
add_image_safe(s1, 'cotomo-logo.png', Inches(4.8), Inches(1.8), height=Inches(0.9))
add_text_box(s1, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2),
             'Cotomo 1分ピッチ', font_size=44, bold=True, color=TEXT_PRIMARY)
add_text_box(s1, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
             'Starley株式会社 ｜ VPoE 篠原 祐貴', font_size=22, color=TEXT_SECONDARY)
# 左下にStarleyロゴ
add_image_safe(s1, 'starley-logo.png', Inches(0.4), Inches(6.8), height=Inches(0.3))

# =================================================================
# S2: プロフィール
# =================================================================
s2 = prs.slides.add_slide(blank_layout)
set_bg_color(s2, BG_LIGHT)
add_bubble(s2, Inches(10.5), Inches(-0.5), Inches(2.5), Inches(2.5), BUBBLE_GREEN)
add_bubble(s2, Inches(-0.3), Inches(5.5), Inches(1.8), Inches(1.8), BUBBLE_PINK)

add_text_box(s2, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             '篠原 祐貴', font_size=40, bold=True, color=TEXT_PRIMARY,
             alignment=PP_ALIGN.LEFT)
add_text_box(s2, Inches(0.8), Inches(1.3), Inches(6), Inches(0.5),
             'Starley株式会社 VPoE', font_size=22, bold=True, color=BRAND_GREEN,
             alignment=PP_ALIGN.LEFT)

bio1 = 'ヤフーを経て、Schoo初代CTO、メドピア技術部長、マネーフォワードケッサイ取締役CTOを歴任。その後SpirのCRO兼CTO、LegalscapeのVPoEを務めるなど、一貫して技術組織の牽引とバリューの最大化に従事。'
bio2 = '現在はStarleyのVPoEとして、エンジニアリングが創出する価値の全体設計を担当。AIと人が共創する「AIネイティブ」なプロダクト組織の実現をリードしています。'
add_text_box(s2, Inches(0.8), Inches(2.0), Inches(10.5), Inches(1.2),
             bio1, font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.LEFT)
add_text_box(s2, Inches(0.8), Inches(3.4), Inches(10.5), Inches(1.2),
             bio2, font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.LEFT)

# キャリアタグ
careers = ['Yahoo!', 'Schoo CTO', 'メドピア', 'MFケッサイ CTO', 'Spir CRO/CTO', 'Legalscape VPoE', 'Starley VPoE']
tag_x = Inches(0.8)
for i, career in enumerate(careers):
    is_current = (career == 'Starley VPoE')
    tag_color = BRAND_GREEN if is_current else TEXT_PRIMARY
    bg_shape = add_rounded_rect(s2, tag_x, Inches(5.0), Inches(1.6), Inches(0.45),
                                 fill_color=RGBColor(0xF0, 0xF0, 0xF0) if not is_current else RGBColor(0xE8, 0xF8, 0xF0))
    t = add_text_box(s2, tag_x + Inches(0.1), Inches(5.05), Inches(1.4), Inches(0.35),
                     career, font_size=13, bold=True, color=tag_color, alignment=PP_ALIGN.CENTER)
    tag_x += Inches(1.7)

# =================================================================
# S3: 会社紹介
# =================================================================
s3 = prs.slides.add_slide(blank_layout)
set_bg_color(s3, BG_LIGHT)
add_bubble(s3, Inches(-0.5), Inches(-0.6), Inches(2.8), Inches(2.8), BUBBLE_YELLOW)
add_bubble(s3, Inches(10.5), Inches(5.5), Inches(2), Inches(2), BUBBLE_GREEN)
add_image_safe(s3, 'starley-logo.png', Inches(4.5), Inches(2.5), height=Inches(0.7))
add_text_box(s3, Inches(1), Inches(4.0), Inches(11.3), Inches(0.7),
             '2023年4月創業 ｜ AI関連プロダクトの企画・開発', font_size=26, color=TEXT_SECONDARY)

# =================================================================
# S4: プロダクト紹介
# =================================================================
s4 = prs.slides.add_slide(blank_layout)
set_bg_color(s4, BG_LIGHT)
add_image_safe(s4, 'bg-characters.png', Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_bubble(s4, Inches(-0.5), Inches(-0.7), Inches(3), Inches(3), BUBBLE_GREEN)
add_bubble(s4, Inches(4), Inches(5.5), Inches(2.2), Inches(2.2), BUBBLE_YELLOW)
add_bubble(s4, Inches(10.5), Inches(0.6), Inches(2), Inches(2), BUBBLE_PINK)

add_text_box(s4, Inches(0.8), Inches(1.5), Inches(5), Inches(2),
             'だれと、\nおしゃべりする？', font_size=38, bold=True, color=TEXT_PRIMARY,
             alignment=PP_ALIGN.LEFT)
add_image_safe(s4, 'cotomo-logo.png', Inches(0.8), Inches(3.8), height=Inches(0.7))
add_text_box(s4, Inches(0.8), Inches(4.8), Inches(5), Inches(0.5),
             'おしゃべりAI コトモ', font_size=24, color=TEXT_SECONDARY,
             alignment=PP_ALIGN.LEFT)
add_image_safe(s4, 'app-mockup.png', Inches(6.5), Inches(0.8), height=Inches(5.5))

# =================================================================
# S5: 機能紹介（2カラム）
# =================================================================
s5 = prs.slides.add_slide(blank_layout)
set_bg_color(s5, BG_LIGHT)

# 左カラムタイトル
add_text_box(s5, Inches(0.5), Inches(0.3), Inches(5.5), Inches(0.8),
             '声、アイコン、性格をカスタマイズして\nAIキャラを作成', font_size=16, bold=True,
             color=RGBColor(0x1B, 0x1B, 0x1B))
# 右カラムタイトル
add_text_box(s5, Inches(7), Inches(0.3), Inches(5.5), Inches(0.8),
             '他のユーザーが作成したキャラと出会う\nプラットフォーム', font_size=16, bold=True,
             color=RGBColor(0x1B, 0x1B, 0x1B))

# グロー装飾
add_image_safe(s5, 'glow-orange.png', Inches(0.3), Inches(1.2), width=Inches(2.5))
add_image_safe(s5, 'glow-blue.png', Inches(3), Inches(4.5), width=Inches(2.5))
add_image_safe(s5, 'glow-pink.png', Inches(8), Inches(1), width=Inches(2.2))
add_image_safe(s5, 'glow-green.png', Inches(7.5), Inches(4), width=Inches(2.2))

# スクリーンショット
add_image_safe(s5, 'feature-settings.png', Inches(1.3), Inches(1.3), height=Inches(5))
add_image_safe(s5, 'feature-characters.png', Inches(7.5), Inches(1.3), height=Inches(5))

# 機能タグ（左カラム）
for tag_data in [
    ('ハイクオリティな', '声', Inches(0.5), Inches(3.0)),
    ('作り込める', '性格', Inches(4.5), Inches(4.5)),
    ('自由に設定', 'アイコン', Inches(4.5), Inches(2.0)),
]:
    sm, lg, tx, ty = tag_data
    add_text_box(s5, tx, ty, Inches(1.5), Inches(0.3), sm, font_size=11, bold=True, color=WHITE)
    add_text_box(s5, tx, ty + Inches(0.25), Inches(1.5), Inches(0.4), lg, font_size=20, bold=True, color=WHITE)

# 機能タグ（右カラム）
add_text_box(s5, Inches(11), Inches(4.5), Inches(1.5), Inches(0.3), 'キャラを', font_size=11, bold=True, color=WHITE)
add_text_box(s5, Inches(11), Inches(4.75), Inches(1.5), Inches(0.4), '公開', font_size=20, bold=True, color=WHITE)

# =================================================================
# S6: ユーザー数
# =================================================================
s6 = prs.slides.add_slide(blank_layout)
set_bg_color(s6, WHITE)
add_rich_text(s6, Inches(1), Inches(2.0), Inches(11.3), Inches(2.5), [
    {'text': '200', 'size': 120, 'bold': True, 'italic': True, 'color': TEXT_PRIMARY},
    {'text': '万', 'size': 64, 'bold': True, 'color': TEXT_PRIMARY},
])
add_text_box(s6, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8),
             'ユーザー突破！', font_size=38, bold=True, color=TEXT_PRIMARY)

# =================================================================
# S7: 累計おしゃべり時間
# =================================================================
s7 = prs.slides.add_slide(blank_layout)
set_bg_color(s7, WHITE)

# キャラクターイラスト
add_image_safe(s7, 'char-boy.png', Inches(0.4), Inches(0.4), height=Inches(2))
add_image_safe(s7, 'char-girl-purple.png', Inches(1.2), Inches(5), height=Inches(1.8))
add_image_safe(s7, 'char-girl-green.png', Inches(10.5), Inches(4.8), height=Inches(1.9))

# 黄色バブル
add_bubble(s7, Inches(3.5), Inches(5), Inches(3), Inches(3), BUBBLE_YELLOW, opacity=0.40)

# タイトル
add_text_box(s7, Inches(2), Inches(1.5), Inches(9.3), Inches(0.6),
             '全ユーザー累計おしゃべり時間', font_size=24, bold=True, color=TEXT_DARK)

# 数値
add_rich_text(s7, Inches(2), Inches(2.2), Inches(9.3), Inches(2.5), [
    {'text': '250', 'size': 110, 'bold': True, 'italic': True, 'color': TEXT_DARK},
    {'text': '万', 'size': 64, 'bold': True, 'color': TEXT_DARK},
    {'text': '時間', 'size': 44, 'bold': True, 'color': TEXT_DARK},
])

# 補足カード
card = add_rounded_rect(s7, Inches(2.5), Inches(4.8), Inches(8.3), Inches(1.1), WHITE)
# カードテキスト（リッチテキスト）
add_rich_text(s7, Inches(2.8), Inches(5.0), Inches(7.8), Inches(0.7), [
    {'text': 'リリースから', 'size': 18, 'color': TEXT_DARK},
    {'text': '1.5年', 'size': 18, 'bold': True, 'color': ACCENT_PINK},
    {'text': 'で、約', 'size': 18, 'color': TEXT_DARK},
    {'text': '285年分', 'size': 18, 'bold': True, 'color': ACCENT_PINK},
    {'text': 'のおしゃべり時間！', 'size': 18, 'color': TEXT_DARK},
])

# =================================================================
# S8: CTA（今すぐおしゃべり）
# =================================================================
s8 = prs.slides.add_slide(blank_layout)
set_bg_color(s8, WHITE)

add_image_safe(s8, 'blob-red.png', Inches(10), Inches(-0.5), width=Inches(3.5))
add_image_safe(s8, 'blob-pink.png', Inches(0.5), Inches(5.5), width=Inches(1.5))

add_image_safe(s8, 'cotomo-logo.png', Inches(5.2), Inches(1.2), height=Inches(0.6))
add_text_box(s8, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
             '今すぐ\nおしゃべり！', font_size=36, bold=True, color=TEXT_DARK)

# 区切り線
line = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(3.8), Inches(1), Pt(1.5))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
line.line.fill.background()

add_text_box(s8, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
             'おしゃべりAI コトモ', font_size=20, color=RGBColor(0x54, 0x54, 0x54))

# QR + URL
add_image_safe(s8, 'qr-code.png', Inches(5.2), Inches(5.0), height=Inches(1.3))
add_text_box(s8, Inches(7), Inches(5.3), Inches(3), Inches(0.4),
             'https://cotomo.ai/', font_size=18, color=TEXT_SECONDARY)

# =================================================================
# S9: 本題タイトル
# =================================================================
s9 = prs.slides.add_slide(blank_layout)
set_bg_color(s9, BG_LIGHT)
add_image_safe(s9, 'bg-characters.png', Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_text_box(s9, Inches(1), Inches(2.8), Inches(11.3), Inches(1.2),
             'Cotomo 1分ピッチ', font_size=44, bold=True, color=TEXT_PRIMARY)

# =================================================================
# S10: 3つだけ覚えてください
# =================================================================
s10 = prs.slides.add_slide(blank_layout)
set_bg_color(s10, BG_LIGHT)
add_bubble(s10, Inches(8.5), Inches(-0.8), Inches(3.3), Inches(3.3), BUBBLE_GREEN)
add_bubble(s10, Inches(-0.4), Inches(5.5), Inches(2.2), Inches(2.2), BUBBLE_YELLOW)
add_bubble(s10, Inches(8), Inches(4.5), Inches(2), Inches(2), BUBBLE_PINK)

add_text_box(s10, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
             '3つだけ、覚えてください。', font_size=44, bold=True, color=TEXT_PRIMARY)

# 3カード
cards = [('Cotomo', BRAND_GREEN), ('AIキャラクター市場', TEXT_PRIMARY), ('採用', TEXT_PRIMARY)]
card_x = Inches(2.5)
for label, color in cards:
    bg = add_rounded_rect(s10, card_x, Inches(3.5), Inches(2.6), Inches(0.8),
                          fill_color=RGBColor(0xFA, 0xFA, 0xFA))
    add_text_box(s10, card_x + Inches(0.1), Inches(3.55), Inches(2.4), Inches(0.7),
                 label, font_size=22, bold=True, color=color)
    card_x += Inches(2.9)

add_note(s10, '最初に、これだけ覚えていてください。Cotomo。AIキャラクターという新しいエンタメ市場。そして、デザイナとエンジニアを本気で募集しています。これだけ覚えてもらえれば大丈夫です。')

# =================================================================
# S11: 聞いたこと、ありますか？
# =================================================================
s11 = prs.slides.add_slide(blank_layout)
set_bg_color(s11, BG_LIGHT)
add_bubble(s11, Inches(-0.5), Inches(-0.6), Inches(2.8), Inches(2.8), BUBBLE_GREEN)
add_bubble(s11, Inches(10.5), Inches(4), Inches(1.8), Inches(1.8), BUBBLE_PINK)

add_image_safe(s11, 'cotomo-logo.png', Inches(5.2), Inches(2.2), height=Inches(0.6))
add_text_box(s11, Inches(1), Inches(3.2), Inches(11.3), Inches(1.2),
             '聞いたこと、ありますか？', font_size=44, bold=True, color=TEXT_PRIMARY)

add_note(s11, 'ひとつだけ聞かせてください。Cotomoという、AIキャラクターとおしゃべりできるアプリを、聞いたことあるかも？という人、どれくらいいますか？（挙手を促す）ありがとうございます。')

# =================================================================
# S12: AIキャラクターという新しいエンタメ市場
# =================================================================
s12 = prs.slides.add_slide(blank_layout)
set_bg_color(s12, BG_LIGHT)
add_bubble(s12, Inches(8.5), Inches(-0.7), Inches(3), Inches(3), BUBBLE_YELLOW)
add_bubble(s12, Inches(0.8), Inches(5.5), Inches(2.2), Inches(2.2), BUBBLE_GREEN)
add_bubble(s12, Inches(-0.3), Inches(1), Inches(1.7), Inches(1.7), BUBBLE_PINK)

# リッチテキスト（「静かに立ち上がっている。」がアクセントカラー）
txbox = s12.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(4))
tf = txbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
for text, color, bold in [
    ('AIキャラクターという\n新しいエンタメ市場が、\n世界的に、', TEXT_PRIMARY, True),
    ('静かに立ち上がっている。', BRAND_GREEN, True),
]:
    r = p.add_run()
    r.text = text
    r.font.size = Pt(40)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = 'Noto Sans JP'

add_note(s12, 'まだ知らない人のほうが多いですよね。でも実は今、AIキャラクターと話し、自分の言葉で物語や体験が変わっていく——そんな新しいエンタメ市場が、世界的に静かに立ち上がっています。')

# =================================================================
# S13: 200万+ DL
# =================================================================
s13 = prs.slides.add_slide(blank_layout)
set_bg_color(s13, BG_LIGHT)
add_bubble(s13, Inches(-0.8), Inches(-1), Inches(3.8), Inches(3.8), BUBBLE_GREEN)
add_bubble(s13, Inches(10), Inches(5.5), Inches(2.2), Inches(2.2), BUBBLE_YELLOW)

add_text_box(s13, Inches(1), Inches(1.8), Inches(11.3), Inches(2),
             '200万+ DL', font_size=110, bold=True, color=BRAND_GREEN)
add_text_box(s13, Inches(1), Inches(4.5), Inches(11.3), Inches(0.7),
             '日本では、今が一番大事なタイミングだ。', font_size=26, color=TEXT_SECONDARY)

add_note(s13, 'リリース直後に100万、200万ダウンロードに到達するプロダクトも出てきていて、日本では、今が一番大事なタイミングです。')

# =================================================================
# S14: 市場の成長と一緒にいるフェーズ
# =================================================================
s14 = prs.slides.add_slide(blank_layout)
set_bg_color(s14, BG_LIGHT)
add_bubble(s14, Inches(9), Inches(-0.6), Inches(2.8), Inches(2.8), BUBBLE_PINK)
add_bubble(s14, Inches(-0.5), Inches(5), Inches(2.5), Inches(2.5), BUBBLE_GREEN)

add_text_box(s14, Inches(1), Inches(1.5), Inches(11.3), Inches(2.5),
             '市場の成長と\n一緒にいるフェーズだ。', font_size=44, bold=True, color=TEXT_PRIMARY)
add_text_box(s14, Inches(1), Inches(4.3), Inches(11.3), Inches(0.8),
             '乗るのではない。作る。', font_size=36, bold=True, color=BRAND_GREEN)

add_note(s14, 'Cotomoは、日本発でこの市場を作りにいっています。成された市場に乗るのではなく、市場の成長と一緒にいるフェーズです。だから正直に言うと——')

# =================================================================
# S15: 一緒に作る側として、助けてください
# =================================================================
s15 = prs.slides.add_slide(blank_layout)
set_bg_color(s15, BG_LIGHT)
add_bubble(s15, Inches(-0.6), Inches(-0.8), Inches(3.3), Inches(3.3), BUBBLE_GREEN)
add_bubble(s15, Inches(9), Inches(5), Inches(2.6), Inches(2.6), BUBBLE_YELLOW)
add_bubble(s15, Inches(8), Inches(0.6), Inches(1.8), Inches(1.8), BUBBLE_PINK)

# リッチテキスト
txbox = s15.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.3), Inches(2.5))
tf = txbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
for text, color in [
    ('一緒に作る側として、\nぜひ', TEXT_PRIMARY),
    ('助けてください。', BRAND_GREEN),
]:
    r = p.add_run()
    r.text = text
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = 'Noto Sans JP'

# CTAバッジ
badge = add_rounded_rect(s15, Inches(4.2), Inches(3.5), Inches(5), Inches(0.7), BRAND_GREEN)
add_text_box(s15, Inches(4.2), Inches(3.55), Inches(5), Inches(0.6),
             'Designer & Engineer 募集', font_size=26, bold=True, color=WHITE)

# QR + Cotomo
add_image_safe(s15, 'qr-code.png', Inches(4.5), Inches(4.8), height=Inches(1.5))
add_image_safe(s15, 'cotomo-logo.png', Inches(6.5), Inches(5.0), height=Inches(0.4))
add_text_box(s15, Inches(6.5), Inches(5.5), Inches(3), Inches(0.4),
             'https://cotomo.ai/', font_size=18, color=TEXT_SECONDARY,
             alignment=PP_ALIGN.LEFT)

add_note(s15, 'フルコミットできるデザイナとエンジニアが、本当に必要です。今しかできないこの挑戦を、一緒に作る側として、ぜひ助けてください。')

# --- 保存 ---
output_path = os.path.join(os.path.dirname(__file__), 'cotomo-pitch.pptx')
prs.save(output_path)
print(f'PPTX生成完了: {output_path}')
print(f'スライド数: {len(prs.slides)}')
